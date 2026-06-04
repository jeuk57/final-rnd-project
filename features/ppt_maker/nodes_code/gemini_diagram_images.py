from __future__ import annotations

import argparse
import base64
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from google import genai
    from google.genai import types
except Exception:  # pragma: no cover - 로컬 미리보기 실행 시 선택 의존성
    genai = None
    types = None
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


IMAGE_PROMPT_BASE = """
public-funded national R&D presentation visual
official technical report tone
institutional and credible style
research program briefing material style
non-decorative and evidence-oriented composition

flat vector infographic
presentation diagram
workflow or process illustration
clear logical flow
simple geometric shapes
clean arrows and connections
balanced layout

professional infographic used in government or research presentation

white or light background
limited color palette (blue / teal / gray)
high contrast icons

NO TEXT
NO LETTERS
NO NUMBERS
All boxes must be empty
Do not render any characters

avoid AI-art look
avoid fantasy style
avoid cinematic rendering
avoid 3D render look
no photorealistic
avoid Korean text
no UI screenshots
""".strip()

# 이미지 생성 프롬프트: Gemini 이미지 모델 호환성을 위해 영어 지시문 유지

BOXES_TOP: List[Tuple[float, float, float, float]] = [
    (0.239746, 0.143072, 0.111328, 0.062444),
    (0.369848, 0.138593, 0.120605, 0.062444),
    (0.500000, 0.142180, 0.120605, 0.062444),
    (0.646854, 0.138593, 0.114746, 0.062444),
    (0.782227, 0.138593, 0.117676, 0.062444),
]

BOXES_PROCESS: List[Tuple[float, float, float, float]] = BOXES_TOP

BOXES_SERVICE: List[Tuple[float, float, float, float]] = [
    (0.252196, 0.291685, 0.085449, 0.052631),
    (0.387329, 0.291685, 0.085449, 0.053520),
    (0.522950, 0.289196, 0.085938, 0.053520),
    (0.664551, 0.288295, 0.084961, 0.053520),
    (0.799684, 0.288295, 0.085938, 0.054421),
]

BOXES_MODEL: List[Tuple[float, float, float, float]] = [
    (0.254150, 0.425799, 0.083984, 0.054421),
    (0.375484, 0.425799, 0.083496, 0.054421),
    (0.480389, 0.426087, 0.089355, 0.053520),
    (0.590591, 0.425591, 0.085449, 0.052631),
    (0.692580, 0.426615, 0.083008, 0.053520),
    (0.801393, 0.426643, 0.082520, 0.053520),
]

BOXES_DATA: List[Tuple[float, float, float, float]] = [
    (0.252196, 0.566899, 0.085938, 0.051740),
    (0.364471, 0.567792, 0.098633, 0.050847),
    (0.492401, 0.566899, 0.077148, 0.051740),
    (0.596191, 0.567792, 0.081055, 0.051740),
    (0.706543, 0.567792, 0.075684, 0.051740),
    (0.806152, 0.567792, 0.084961, 0.052632),
]

BOXES_USERS: List[Tuple[float, float, float, float]] = [
    (0.050000, 0.227980, 0.115000, 0.022000),
    (0.050000, 0.396000, 0.115000, 0.022000),
    (0.050000, 0.564331, 0.115000, 0.022000),
    (0.050000, 0.734000, 0.115000, 0.022000),
]

BOXES_DATA_ICONS: List[Tuple[float, float, float, float]] = [
    (0.236863, 0.727000, 0.088000, 0.018000),
    (0.354363, 0.725000, 0.095000, 0.018000),
    (0.465535, 0.725000, 0.079000, 0.018000),
    (0.590591, 0.727000, 0.082000, 0.018000),
    (0.704227, 0.725000, 0.078000, 0.018000),
    (0.818297, 0.727148, 0.086000, 0.018000),
]

BOXES_INFRA_ICONS: List[Tuple[float, float, float, float]] = [
    (0.280863, 0.935000, 0.083000, 0.035000),
    (0.378177, 0.935000, 0.083000, 0.035000),
    (0.473677, 0.935000, 0.083000, 0.035000),
    (0.567386, 0.935000, 0.083000, 0.035000),
    (0.654304, 0.937069, 0.083000, 0.035000),
    (0.749512, 0.937080, 0.083000, 0.035000),
]

# 시스템 아키텍처는 고정 배경 이미지만 사용하고 라벨 오버레이는 비활성화
LABELS_SERVICE: List[str] = []
LABELS_MODEL: List[str] = []
LABELS_PROCESS: List[str] = []
LABELS_DATA: List[str] = []
LABELS_USERS: List[str] = []
LABELS_DATA_ICONS: List[str] = []
LABELS_INFRA_ICONS: List[str] = []

def _norm(s: Any) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


# 문자열 bool 설정값 해석 함수
def _to_bool(v: Any) -> Optional[bool]:
    if isinstance(v, bool):
        return v
    if v is None:
        return None
    t = str(v).strip().lower()
    if t in {"1", "true", "yes", "y", "on"}:
        return True
    if t in {"0", "false", "no", "n", "off"}:
        return False
    return None


# Gemini 이미지 생성 사용 여부 판단 함수
def _enabled(state: Optional[Dict[str, Any]]) -> bool:
    if isinstance(state, dict):
        v = _to_bool(state.get("enable_gemini_diagram_images"))
        if v is not None:
            return v
    v = _to_bool(os.environ.get("ENABLE_GEMINI_DIAGRAM_IMAGES"))
    if v is not None:
        return v
    return True


# Gemini 응답 객체에서 이미지 바이트 추출 함수
def _extract_image_bytes(resp: Any) -> Optional[bytes]:
    for cand in (getattr(resp, "candidates", None) or []):
        content = getattr(cand, "content", None)
        for part in (getattr(content, "parts", None) or []):
            inline = getattr(part, "inline_data", None)
            if not inline:
                continue
            data = getattr(inline, "data", None)
            if not data:
                continue
            if isinstance(data, (bytes, bytearray)):
                return bytes(data)
            if isinstance(data, str):
                try:
                    return base64.b64decode(data)
                except Exception:
                    continue
    return None


# 이미지 생성 가능 Gemini 모델 후보 탐색 함수
def _discover_model_candidates(client: genai.Client, preferred: str) -> List[str]:
    raw: List[str] = [preferred]
    try:
        for m in list(client.models.list()):
            name = str(getattr(m, "name", "") or "").strip()
            if not name:
                continue
            actions = list(getattr(m, "supported_actions", None) or [])
            if "generateContent" not in actions:
                continue
            lk = name.lower()
            if ("image" in lk) or ("imagen" in lk) or ("flash-exp-image-generation" in lk):
                raw.append(name)
    except Exception as e:
        print(f"[WARN] 모델 목록 조회 실패, 고정 후보 사용: {e}")

    raw.extend(
        [
            "models/gemini-2.5-flash-image",
            "models/gemini-2.0-flash-exp-image-generation",
            "models/gemini-2.5-flash-image-preview",
            "models/gemini-2.0-flash-preview-image-generation",
            "gemini-2.5-flash-image",
            "gemini-2.0-flash-exp-image-generation",
        ]
    )
    out: List[str] = []
    for m in raw:
        m = (m or "").strip()
        if m and m not in out:
            out.append(m)
    return out


# Gemini 이미지 생성 요청 함수
def _try_generate_with_config(client: genai.Client, model: str, prompt: str, mode: str) -> Optional[bytes]:
    if mode == "IMAGE_ONLY":
        resp = client.models.generate_content(
            model=model,
            contents=prompt,
            config=types.GenerateContentConfig(response_modalities=["IMAGE"], temperature=0.2),
        )
    else:
        return None
    return _extract_image_bytes(resp)


# 단일 이미지 생성 및 저장 함수
def _generate_one_image(
    client: genai.Client,
    models: List[str],
    prompt: str,
    out_path: Path,
    *,
    max_retries: int = 1,
) -> Optional[str]:
    for model in models:
        for mode in ["IMAGE_ONLY"]:
            for attempt in range(max_retries):
                try:
                    img = _try_generate_with_config(client, model, prompt, mode)
                    if not img:
                        continue
                    out_path.parent.mkdir(parents=True, exist_ok=True)
                    out_path.write_bytes(img)
                    print(f"[INFO] Gemini image model selected: {model} ({mode})")
                    return str(out_path)
                except Exception as e:
                    print(f"[WARN] Gemini image generation failed: {model} ({mode}) -> {e}")
                    continue
    return None



# 활용방안/기대효과 이미지 대상 슬라이드 탐색 함수
def _find_effect_slide_idx(deck_slides: List[Dict[str, Any]]) -> Optional[int]:
    for idx, spec in enumerate(deck_slides):
        if not isinstance(spec, dict):
            continue
        section = _norm(spec.get("section"))
        title = _norm(spec.get("slide_title"))
        merged = f"{section} {title}"
        if "활용방안 및 기대효과" in merged or "기대효과" in merged:
            return idx
    return None


# 연구 개요 마지막 이미지 대상 슬라이드 탐색 함수
def _find_overview_slide_idx(deck_slides: List[Dict[str, Any]]) -> Optional[int]:
    for idx, spec in enumerate(deck_slides):
        if not isinstance(spec, dict):
            continue
        if str(spec.get("image_prompt_type") or "") == "overview_last":
            return idx
    return None


# 추진 계획 조직도 대상 슬라이드 탐색 함수
def _find_plan_orgchart_slide_idx(deck_slides: List[Dict[str, Any]]) -> Optional[int]:
    for idx, spec in enumerate(deck_slides):
        if not isinstance(spec, dict):
            continue
        if str(spec.get("image_prompt_type") or "") == "plan_orgchart_fixed":
            return idx
    return None


# 시스템 아키텍처 대상 슬라이드 탐색 함수
def _find_system_arch_slide_idx(deck_slides: List[Dict[str, Any]]) -> Optional[int]:
    for idx, spec in enumerate(deck_slides):
        if not isinstance(spec, dict):
            continue
        if str(spec.get("image_prompt_type") or "") == "system_architecture":
            return idx
        if "시스템 아키텍처" in _norm(spec.get("slide_title")):
            return idx
    return None


# Gemini 이미지 생성 프롬프트 구성 함수
def _build_prompt(deck_title: str, section: str, title: str, prompt_type: str = "") -> str:
    if prompt_type == "system_architecture":
        return (
            "Create a complex system architecture diagram for a Korean government R&D evaluation presentation.\n"
            "Style: clean 2D vector infographic, institutional/technical report tone, minimal, no decorative art, no 3D, "
            "no cinematic lighting, no heavy gradients, no AI-art look.\n"
            "Background: white or very light gray with generous margins, crisp borders, perfectly aligned.\n"
            "Layout (three-tier, like a portal/system blueprint):\n"
            "Left column: four user groups represented only by icons (agency/org, experts, partner institutions, public users). "
            "From each group, draw 1-2 arrows pointing to the central platform.\n"
            "Center: one large frame containing 3-4 section panels (authentication/roles, service modules, data/content repositories, "
            "unified search/analytics/management). Inside each panel, place 6-12 small rounded rectangles as modules, densely but neatly arranged, ALL EMPTY (no labels).\n"
            "Repositories panel: include 6-10 database cylinder icons, NO labels.\n"
            "Bottom row: 6-8 infrastructure icons aligned horizontally (app server, search, community, mail, integration DB, content storage, etc.) "
            "connected with thin lines to the center.\n"
            "Visual rules: consistent stroke width, right-angle connectors, simple geometric icons, complex yet organized appearance.\n"
            "Text rule: NO TEXT anywhere (no English, no Korean, no acronyms). Empty boxes only. "
            "If any text appears, it is a failure.\n"
            "Use pictograms or icons to indicate meaning instead of labels.\n"
            "Output: 16:9 slide-ready, high resolution, sharp vector look.\n"
            "Negative prompt:\n"
            "text, letters, numbers, words, labels, captions, acronyms, Korean text, watermark, photorealistic, 3D, cinematic, anime, cartoon, fantasy"
        )
    if prompt_type == "plan":
        context = " ".join([_norm(deck_title), _norm(section), _norm(title)]).strip()
        return (
            f"Create one workflow/process concept image for a presentation slide context: {context}. "
            "public-funded national R&D presentation visual\n"
            "official technical report tone\n"
            "institutional and credible style\n"
            "research program briefing material style\n"
            "non-decorative and evidence-oriented composition\n\n"
            "flat vector infographic\n"
            "presentation diagram\n"
            "workflow or process illustration\n"
            "clear logical flow\n"
            "simple geometric shapes\n"
            "clean arrows and connections\n"
            "balanced layout\n\n"
            "professional infographic used in government or research presentation\n\n"
            "white or light background\n"
            "limited color palette (blue / teal / gray)\n"
            "high contrast icons\n\n"
            "avoid AI-art look\n"
            "avoid fantasy style\n"
            "avoid cinematic rendering\n"
            "avoid 3D render look\n"
            "no photorealistic\n"
            "avoid Korean text\n"
            "no UI screenshots\n"
            "NO TEXT. Use icons and shapes only. Empty panels only.\n"
            "Use pictograms or icons to indicate meaning instead of labels."
        )
    if prompt_type == "overview_last":
        context = " ".join([_norm(deck_title), _norm(section), _norm(title)]).strip()
        return (
            f"Create one concept image for a presentation slide context: {context}. "
            "public-funded national R&D presentation visual\n"
            "official technical report tone\n"
            "clean 2D vector infographic style\n"
            "balanced composition with central focus\n"
            "Fill most of the canvas with content; avoid large empty margins.\n"
            "MUST use pure white background (#FFFFFF)\n"
            "no dark background, no black background, no gradient background\n"
            "limited palette (blue/teal/gray)\n"
            "English text only if absolutely necessary.\n"
            "Prefer NO TEXT.\n"
            "If text is used, at most 2 labels total, each 1-2 English words.\n"
            "Never use Korean text.\n"
            "avoid long sentences\n"
            "avoid photorealistic style\n"
            "avoid 3D and cinematic look\n"
            "slide-ready 16:9 composition"
        )

    context = " ".join([_norm(deck_title), _norm(section), _norm(title)]).strip()
    return (
        f"Create one concept image for a presentation slide context: {context}. "
        f"{IMAGE_PROMPT_BASE}. "
        "NO TEXT. Use icons and shapes only."
    )


# 일반 텍스트-이미지 레이아웃 이미지 영역 계산 함수
def _text_image_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    left = int(slide_w * 0.60)
    top = int(slide_h * 0.18)
    width = int(slide_w * 0.36)
    height = int(slide_h * 0.64)
    return left, top, width, height


# 시스템 아키텍처 이미지 영역 계산 함수
def _arch_image_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    # 상단 제목 영역 + 하단 전체 이미지 패널
    left = int(slide_w * 0.06)
    top = int(slide_h * 0.18)
    width = int(slide_w * 0.88)
    height = int(slide_h * 0.74)
    return left, top, width, height


# 추진 계획 이미지 영역 계산 함수
def _plan_image_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    # 상단 제목/요약 2줄 + 하단 큰 이미지 패널
    left = int(slide_w * 0.06)
    top = int(slide_h * 0.23)
    width = int(slide_w * 0.88)
    height = int(slide_h * 0.72)  # 약 70~80% 시각 영역
    return left, top, width, height


# 전체 슬라이드 영역 계산 함수
def _full_slide_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    return 0, 0, int(slide_w), int(slide_h)


# basic 배경의 메인 패널 영역 계산 함수
def _basic_main_panel_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    # features/ppt_maker/background/basic_origin.png 기준 콘텐츠 패널 영역
    return (
        int(slide_w * 0.034765625),
        int(slide_h * 0.05625),
        int(slide_w * 0.93046875),
        int(slide_h * 0.8875),
    )


# 표지 이미지 영역 계산 함수
def _cover_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    # 표지 제목 영역 보호를 위한 우측 이미지 영역
    left = int(slide_w * 0.58)
    top = int(slide_h * 0.22)
    width = int(slide_w * 0.36)
    height = int(slide_h * 0.58)
    return left, top, width, height


# 이미지 삽입 후 뒤쪽 배치 함수
def _insert_picture(slide, image_path: str, slot: Tuple[int, int, int, int]) -> bool:
    try:
        left, top, width, height = slot
        pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        # 텍스트가 이미지 위에 오도록 계층 순서 유지
        sp = pic._element
        parent = sp.getparent()
        parent.remove(sp)
        parent.insert(2, sp)
        return True
    except Exception as e:
        print(f"[WARN] failed to place image: {e}")
        return False


# 이미지 삽입 후 앞쪽 배치 함수
def _insert_picture_front(slide, image_path: str, slot: Tuple[int, int, int, int]) -> bool:
    try:
        left, top, width, height = slot
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        return True
    except Exception as e:
        print(f"[WARN] failed to place image(front): {e}")
        return False


# 비율 좌표를 EMU 좌표로 변환하는 함수
def _ratio_to_emu(slide_w: int, slide_h: int, box: Tuple[float, float, float, float]) -> Tuple[int, int, int, int]:
    x, y, w, h = box
    return int(slide_w * x), int(slide_h * y), int(slide_w * w), int(slide_h * h)


# 슬롯 내부 비율 좌표를 EMU 좌표로 변환하는 함수
def _ratio_to_emu_in_slot(
    slot: Tuple[int, int, int, int], box: Tuple[float, float, float, float]
) -> Tuple[int, int, int, int]:
    sx, sy, sw, sh = slot
    x, y, w, h = box
    return sx + int(sw * x), sy + int(sh * y), int(sw * w), int(sh * h)


# 기존 오버레이 텍스트박스 제거 함수
def _remove_overlay_textboxes(slide, prefix: str = "__Overlay__") -> None:
    for sh in list(slide.shapes):
        name = str(getattr(sh, "name", "") or "")
        if not name.startswith(prefix):
            continue
        try:
            sh._element.getparent().remove(sh._element)
        except Exception:
            continue


# 문단 전체 폰트 적용 함수
def _apply_font_all_runs(paragraph, *, size_pt: int, bold: bool, name: str = "Malgun Gothic") -> None:
    runs = list(paragraph.runs)
    if not runs:
        paragraph.font.name = name
        paragraph.font.size = Pt(size_pt)
        paragraph.font.bold = bool(bold)
        return
    for r in runs:
        r.font.name = name
        r.font.size = Pt(size_pt)
        r.font.bold = bool(bold)


# 절대 좌표 기반 라벨 오버레이 함수
def overlay_labels(
    slide,
    slide_w: int,
    slide_h: int,
    boxes: List[Tuple[float, float, float, float]],
    labels: List[str],
    font_size: int,
    *,
    debug: bool = False,
    bold: bool = False,
    shape_name_prefix: str = "__Overlay__",
) -> None:
    def _normalize_label_text(text: Any) -> str:
        # Some upstream strings contain vertical-tab (\x0b) for manual line breaks.
        raw_text = str(text or "").replace("\x0b", "\n")
        # 수동 줄바꿈 유지로 2줄 라벨 강제
        cleaned = re.sub(r"[\x00-\x09\x0c-\x1f]", " ", raw_text)
        lines = [re.sub(r"\s+", " ", str(ln or "")).strip() for ln in cleaned.splitlines()]
        lines = [ln for ln in lines if ln]
        return "\n".join(lines) if lines else ""

    def _fit_font(base_size: int, text: str) -> int:
        # 시스템 아키텍처 라벨 크기 고정으로 갑작스러운 크기 변화 방지
        return int(base_size)

    for i, (box, raw) in enumerate(zip(boxes, labels)):
        left, top, width, height = _ratio_to_emu(slide_w, slide_h, box)
        shp = slide.shapes.add_textbox(left, top, width, height)
        try:
            shp.name = f"{shape_name_prefix}label_{i}"
        except Exception:
            pass
        shp.fill.background()
        if debug:
            shp.line.fill.solid()
            shp.line.fill.fore_color.rgb = RGBColor(255, 0, 0)
            shp.line.width = Pt(1)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = _normalize_label_text(raw)
        _apply_font_all_runs(p, size_pt=_fit_font(font_size, raw), bold=bool(bold))


# 슬롯 내부 좌표 기반 라벨 오버레이 함수
def overlay_labels_in_slot(
    slide,
    slot: Tuple[int, int, int, int],
    boxes: List[Tuple[float, float, float, float]],
    labels: List[str],
    font_size: int,
    *,
    debug: bool = False,
    bold: bool = False,
    shape_name_prefix: str = "__Overlay__",
) -> None:
    def _normalize_label_text(text: Any) -> str:
        raw_text = str(text or "").replace("\x0b", "\n")
        cleaned = re.sub(r"[\x00-\x09\x0c-\x1f]", " ", raw_text)
        lines = [re.sub(r"\s+", " ", str(ln or "")).strip() for ln in cleaned.splitlines()]
        lines = [ln for ln in lines if ln]
        return "\n".join(lines) if lines else ""

    def _fit_font(base_size: int, text: str) -> int:
        return int(base_size)

    for i, (box, raw) in enumerate(zip(boxes, labels)):
        left, top, width, height = _ratio_to_emu_in_slot(slot, box)
        shp = slide.shapes.add_textbox(left, top, width, height)
        try:
            shp.name = f"{shape_name_prefix}label_{i}"
        except Exception:
            pass
        shp.fill.background()
        if debug:
            shp.line.fill.solid()
            shp.line.fill.fore_color.rgb = RGBColor(255, 0, 0)
            shp.line.width = Pt(1)
        else:
            shp.line.fill.background()
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = _normalize_label_text(raw)
        _apply_font_all_runs(p, size_pt=_fit_font(font_size, raw), bold=bool(bold))


PLAN_ORGCHART_TEXT_SPECS: List[Dict[str, Any]] = [
    {"box": (0.117982, 0.18444, 0.45, 0.13), "text": "\uc5f0\uad6c\uac1c\ubc1c\uacfc\uc81c\n\ucc28\uc138\ub300 \ud574\uc591\u00b7\uadf9\uc9c0 \uae30\ud6c4\uc608\uce21\uc2dc\uc2a4\ud15c \uac1c\ubc1c", "size": 20},
    {"box": (0.652018, 0.2, 0.23, 0.11), "text": "\ucd1d \ucc38\uc5ec \uc5f0\uad6c\uc6d0\n\ucd1d 52\uba85", "size": 20},
    {"box": (0.374141, 0.427226, 0.36, 0.12), "text": "\ud55c\uad6d\ud574\uc591\uacfc\ud559\uae30\uc220\uc6d0(KIOST)\n\ucd1d\uad04\u00b7\uc804\uc9c0\uad6c \ubaa8\ub378\u00b7\ud1b5\ud569 \ud50c\ub7ab\ud3fc", "size": 20},
    {"box": (0.130529, 0.694575, 0.112784, 0.13), "text": "\uc11c\uc6b8\ub300\ud559\uad50\n\uc9c0\uc5ed \uc0c1\uc138\n\ubaa8\ub378\u00b7\uacb0\ud569 \uae30\uc220", "size": 16},
    {"box": (0.321656, 0.692089, 0.2, 0.13), "text": "\uadf9\uc9c0\uc5f0\uad6c\uc18c\n\uadf9\uc9c0 \uc608\uce21\u00b7\ud604\uc7a5 \uad00\uce21", "size": 16},
    {"box": (0.5391, 0.696188, 0.1218, 0.13), "text": "\u321c\ud574\uc591\uc778\uacf5\uc9c0\ub2a5\n\ud574\uc591\uc0dd\ud0dc\uacc4\n\ubaa8\ub378\u00b7\uc5f0\uacc4 \ubd84\uc11d", "size": 16},
    {"box": (0.73045, 0.694575, 0.151568, 0.13), "text": "\u321c\uae30\ud6c4\uc18c\ud504\ud2b8\uc6e8\uc5b4\n\uc6b4\uc601 \ud50c\ub7ab\ud3fc\u00b7\uac00\uc2dc\ud654\n\uc778\ud130\ud398\uc774\uc2a4", "size": 16},
]


# 추진 계획 조직도 라벨 오버레이 함수
def _overlay_plan_orgchart_texts(
    slide,
    slot: Tuple[int, int, int, int],
    specs: Optional[List[Dict[str, Any]]] = None,
) -> None:
    entries = specs or PLAN_ORGCHART_TEXT_SPECS
    _remove_overlay_textboxes(slide, prefix="__OverlayPlan__")
    for item in entries:
        box = tuple(item.get("box") or (0.0, 0.0, 0.0, 0.0))
        text = str(item.get("text") or "").strip()
        if not text:
            continue
        size = int(item.get("size") or 12)
        left, top, width, height = _ratio_to_emu_in_slot(slot, box)  # type: ignore[arg-type]
        shp = slide.shapes.add_textbox(left, top, width, height)
        try:
            shp.name = "__OverlayPlan__label"
        except Exception:
            pass
        shp.fill.background()
        shp.line.fill.background()
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        lines = [re.sub(r"\s+", " ", ln).strip() for ln in text.splitlines() if re.sub(r"\s+", " ", ln).strip()]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.text = ln
            _apply_font_all_runs(p, size_pt=(size if i == 0 else max(10, size - 2)), bold=bool(i == 0))


# 시스템 아키텍처 슬라이드의 제목 외 요소 정리 함수
def _clear_arch_slide_except_title(slide, slide_w: int, slide_h: int):
    _keep_only_title_text(slide, slide_w=slide_w, slide_h=slide_h)
    title_shape = None
    for sh in list(slide.shapes):
        if getattr(sh, "has_text_frame", False):
            title_shape = sh
            break
    # 시스템 아키텍처 배경의 좌측 아이콘 영역과 제목 겹침 방지
    if title_shape is not None:
        try:
            raw_title = _norm(getattr(title_shape, "text", ""))
            if len(raw_title) > 24:
                title_shape.text_frame.clear()
                title_shape.text_frame.paragraphs[0].text = "시스템 아키텍처"
            title_shape.left = int(slide_w * 0.03)
            title_shape.top = int(slide_h * 0.02)
            title_shape.width = int(slide_w * 0.24)
            title_shape.height = int(slide_h * 0.08)
        except Exception:
            pass
    for sh in list(slide.shapes):
        if title_shape is not None and sh == title_shape:
            continue
        try:
            el = sh._element
            el.getparent().remove(el)
        except Exception:
            continue
    return title_shape


# 시스템 아키텍처 배경 이미지 경로 해석 함수
def _resolve_arch_bg_image_path(img_path: str, state: Optional[Dict[str, Any]] = None) -> str:
    candidates: List[Path] = []
    env_path = str(os.environ.get("SYSTEM_ARCH_BG_IMAGE_PATH", "")).strip()
    state_path = str((state or {}).get("system_arch_bg_image_path") or "").strip()
    if state_path:
        candidates.append(Path(state_path))
    if env_path:
        candidates.append(Path(env_path))
    if img_path:
        candidates.append(Path(img_path))
    candidates.append(Path("features/ppt_maker/background/system.png"))
    candidates.append(Path("output/images/system_architecture_bg.png"))
    for c in candidates:
        if c.exists():
            return str(c)
    return str(Path(img_path)) if img_path else str(Path("features/ppt_maker/background/system.png"))


# 추진 계획 조직도 배경 이미지 경로 해석 함수
def _resolve_plan_orgchart_bg_image_path(img_path: str, state: Optional[Dict[str, Any]] = None) -> str:
    candidates: List[Path] = []
    env_path = str(os.environ.get("PLAN_ORGCHART_BG_IMAGE_PATH", "")).strip()
    state_path = str((state or {}).get("plan_orgchart_bg_image_path") or "").strip()
    if state_path:
        candidates.append(Path(state_path))
    if env_path:
        candidates.append(Path(env_path))
    if img_path:
        candidates.append(Path(img_path))
    candidates.append(Path("output/images/plan_orgchart_bg.png"))
    candidates.append(Path("features/ppt_maker/background/Hierarchical organization diagram..png"))
    for c in candidates:
        if c.exists():
            return str(c)
    return str(Path(img_path)) if img_path else str(Path("output/images/plan_orgchart_bg.png"))


# basic 계열 배경 이미지 경로 해석 함수
def _resolve_basic_origin_bg_path(state: Optional[Dict[str, Any]] = None) -> str:
    candidates: List[Path] = []
    state_path = str((state or {}).get("postprocess_bg_basic_image") or "").strip()
    env_path = str(os.environ.get("POSTPROCESS_BG_BASIC_IMAGE") or "").strip()
    legacy_state_path = str((state or {}).get("postprocess_bg_basic_origin") or "").strip()
    legacy_env_path = str(os.environ.get("POSTPROCESS_BG_BASIC_ORIGIN") or "").strip()
    if state_path:
        candidates.append(Path(state_path))
    if env_path:
        candidates.append(Path(env_path))
    if legacy_state_path:
        candidates.append(Path(legacy_state_path))
    if legacy_env_path:
        candidates.append(Path(legacy_env_path))
    candidates.append(Path("features/ppt_maker/background/basic_image.png"))
    candidates.append(Path("features/ppt_maker/background/basic_origin.png"))
    for c in candidates:
        if c.exists():
            return str(c)
    return str(Path("features/ppt_maker/background/basic_image.png"))


# basic 배경 적용 함수
def _apply_basic_origin_background(slide, slide_w: int, slide_h: int, state: Optional[Dict[str, Any]] = None) -> bool:
    bg_path = _resolve_basic_origin_bg_path(state=state)
    if not Path(bg_path).exists():
        return False
    _remove_randi_background_shapes(slide)
    try:
        pic = slide.shapes.add_picture(bg_path, 0, 0, width=int(slide_w), height=int(slide_h))
        try:
            pic.name = "__RandiBgImage__"
        except Exception:
            pass
        sp = pic._element
        parent = sp.getparent()
        parent.remove(sp)
        parent.insert(2, sp)
        return True
    except Exception:
        return False


# 시스템 아키텍처 고정 이미지 및 라벨 오버레이 적용 함수
def add_architecture_overlay(
    slide,
    img_path: str,
    state: Optional[Dict[str, Any]] = None,
    slide_w: Optional[int] = None,
    slide_h: Optional[int] = None,
) -> bool:
    try:
        if slide_w is None or slide_h is None:
            raise RuntimeError("시스템 아키텍처 오버레이에는 slide_w/slide_h가 필요합니다.")
        slide_w = int(slide_w)
        slide_h = int(slide_h)
        _keep_only_title_text(slide, slide_w=slide_w, slide_h=slide_h)
        _apply_basic_origin_background(slide, slide_w, slide_h, state=state)
        slot = _basic_main_panel_slot(slide_w, slide_h)
        _remove_overlay_textboxes(slide, prefix="__OverlayArch__")
        _remove_overlapping_shapes(
            slide,
            slot,
            slide_w=slide_w,
            slide_h=slide_h,
            margin_ratio=0.0,
            preserve_text_shapes=False,
        )
        bg_path = _resolve_arch_bg_image_path(img_path, state=state)
        if not _insert_picture_front(slide, bg_path, slot):
            return False

        debug = _to_bool((state or {}).get("arch_overlay_debug"))
        if debug is None:
            debug = bool(_to_bool(os.environ.get("ARCH_OVERLAY_DEBUG")))
        # 라벨 폰트 정책: 박스 내부 12pt, 외부 보조 라벨 9pt, 전체 굵게
        overlay_labels_in_slot(slide, slot, BOXES_TOP, LABELS_PROCESS, 12, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__top_")
        overlay_labels_in_slot(slide, slot, BOXES_SERVICE, LABELS_SERVICE, 12, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__service_")
        overlay_labels_in_slot(slide, slot, BOXES_MODEL, LABELS_MODEL, 12, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__model_")
        overlay_labels_in_slot(slide, slot, BOXES_DATA, LABELS_DATA, 12, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__data_")
        overlay_labels_in_slot(slide, slot, BOXES_USERS, LABELS_USERS, 9, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__users_")
        overlay_labels_in_slot(slide, slot, BOXES_DATA_ICONS, LABELS_DATA_ICONS, 9, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__dataicons_")
        overlay_labels_in_slot(slide, slot, BOXES_INFRA_ICONS, LABELS_INFRA_ICONS, 9, debug=bool(debug), bold=True, shape_name_prefix="__OverlayArch__infra_")
        return True
    except Exception as e:
        print(f"[WARN] 시스템 아키텍처 오버레이 실패, 기존 이미지 삽입 흐름으로 대체: {e}")
        return False


# 이미지 영역과 겹치는 기존 도형 제거 함수
def _remove_overlapping_shapes(
    slide,
    slot: Tuple[int, int, int, int],
    *,
    slide_w: int,
    slide_h: int,
    margin_ratio: float = 0.06,
    preserve_text_shapes: bool = False,
) -> None:
    # 이미지 영역 주변 6% 여백까지 제거 대상으로 포함해 가장자리 겹침 방지
    mx = int(slide_w * margin_ratio)
    my = int(slide_h * margin_ratio)
    left, top, width, height = slot
    left = max(0, left - mx)
    top = max(0, top - my)
    right = min(slide_w, left + width + (mx * 2))
    bottom = min(slide_h, top + height + (my * 2))

    def _intersects(sh) -> bool:
        try:
            l2, t2 = int(sh.left), int(sh.top)
            r2, b2 = l2 + int(sh.width), t2 + int(sh.height)
        except Exception:
            return False
        return max(left, l2) < min(right, r2) and max(top, t2) < min(bottom, b2)

    for sh in list(slide.shapes):
        name = str(getattr(sh, "name", "") or "")
        if name in {"__RandiBgImage__", "RandiBackground"}:
            continue
        if not _intersects(sh):
            continue
        if preserve_text_shapes and getattr(sh, "has_text_frame", False):
            continue
        # 임시 영역 처리:
        # - 텍스트 임시 영역은 텍스트만 비우고 보존 가능
        # - 그림/객체 임시 영역은 회색 박스 방지를 위해 제거
        if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PLACEHOLDER:
            phf = getattr(sh, "placeholder_format", None)
            pht = int(getattr(phf, "type", -1)) if phf is not None else -1
            # 텍스트 placeholder 유형: TITLE(1), CENTER_TITLE(3), SUBTITLE(4), BODY(2), VERTICAL 계열(5~8)
            text_placeholder_types = {1, 2, 3, 4, 5, 6, 7, 8}
            if pht in text_placeholder_types and getattr(sh, "has_text_frame", False):
                try:
                    sh.text_frame.clear()
                except Exception:
                    pass
                continue
        # 겹치는 도형 전체 제거(그룹 도형 포함)
        try:
            el = sh._element
            el.getparent().remove(el)
        except Exception:
            try:
                # 복합 도형 참조 제거 보완
                slide.shapes._spTree.remove(sh._element)  # pylint: disable=protected-access
            except Exception:
                continue


# 기존 배경 이미지 도형 제거 함수
def _remove_randi_background_shapes(slide) -> int:
    removed = 0
    for sh in list(slide.shapes):
        name = str(getattr(sh, "name", "") or "")
        if name not in {"__RandiBgImage__", "RandiBackground"}:
            continue
        try:
            el = sh._element
            el.getparent().remove(el)
            removed += 1
        except Exception:
            continue
    return removed


# 표지 이미지 영역과 겹치는 보조 텍스트 정리 함수
def _trim_cover_secondary_text(slide, slot: Tuple[int, int, int, int], *, slide_w: int, slide_h: int) -> None:
    left, top, width, height = slot
    right = left + width
    bottom = top + height

    def _intersects(sh) -> bool:
        try:
            l2, t2 = int(sh.left), int(sh.top)
            r2, b2 = l2 + int(sh.width), t2 + int(sh.height)
        except Exception:
            return False
        return max(left, l2) < min(right, r2) and max(top, t2) < min(bottom, b2)

    text_shapes = [sh for sh in list(slide.shapes) if getattr(sh, "has_text_frame", False)]
    if not text_shapes:
        return

    # 좌측 상단 제목성 텍스트는 보존하고 이미지 영역을 침범한 보조 텍스트만 제거
    for sh in text_shapes:
        try:
            l2, t2 = int(sh.left), int(sh.top)
            r2 = l2 + int(sh.width)
            in_title_safe = (l2 < int(slide_w * 0.52)) and (t2 < int(slide_h * 0.40))
        except Exception:
            continue
        if in_title_safe:
            continue
        if _intersects(sh):
            try:
                el = sh._element
                el.getparent().remove(el)
            except Exception:
                try:
                    sh.text_frame.clear()
                except Exception:
                    pass


# 제목 텍스트만 남기고 나머지 텍스트 제거 함수
def _keep_only_title_text(slide, *, slide_w: int, slide_h: int) -> None:
    text_shapes = [sh for sh in list(slide.shapes) if getattr(sh, "has_text_frame", False)]
    if not text_shapes:
        return
    text_shapes.sort(key=lambda sh: (int(getattr(sh, "top", 0)), int(getattr(sh, "left", 0))))
    title_shape = text_shapes[0]
    try:
        title_shape.left = int(slide_w * 0.06)
        title_shape.top = int(slide_h * 0.05)
        title_shape.width = int(slide_w * 0.88)
        title_shape.height = int(slide_h * 0.10)
    except Exception:
        pass
    for sh in text_shapes[1:]:
        try:
            el = sh._element
            el.getparent().remove(el)
        except Exception:
            try:
                sh.text_frame.clear()
            except Exception:
                pass


# 추진 계획 슬라이드 상단 제목/요약 렌더링 함수
def _render_plan_text_header(slide, *, title: str, line1: str, line2: str, slide_w: int, slide_h: int) -> None:
    # Remove existing text shapes and rebuild compact top header.
    for sh in list(slide.shapes):
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            el = sh._element
            el.getparent().remove(el)
        except Exception:
            continue

    title_box = slide.shapes.add_textbox(
        int(slide_w * 0.06),
        int(slide_h * 0.05),
        int(slide_w * 0.88),
        int(slide_h * 0.10),
    )
    title_box.text_frame.clear()
    title_box.text_frame.paragraphs[0].text = _norm(title) or "추진 계획"

    sub_box = slide.shapes.add_textbox(
        int(slide_w * 0.06),
        int(slide_h * 0.14),
        int(slide_w * 0.88),
        int(slide_h * 0.08),
    )
    tf = sub_box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = _norm(line1)
    p2 = tf.add_paragraph()
    p2.text = _norm(line2)


# 연구 개요 중앙 이미지 영역 계산 함수
def _overview_center_image_slot(slide_w: int, slide_h: int) -> Tuple[int, int, int, int]:
    # Centered large visual block for "연구 개요" image slide.
    w = int(slide_w * 0.70)
    h = int(slide_h * 0.60)
    l = int((slide_w - w) * 0.5)
    t = int(slide_h * 0.18)
    return l, t, w, h


# 연구 개요 하단 텍스트 그룹 렌더링 함수
def _overview_bottom_text_groups(slide, spec: Dict[str, Any], *, slide_w: int, slide_h: int) -> None:
    # Build two compact groups under the centered image.
    bullets = [_norm(x) for x in (spec.get("bullets") or []) if _norm(x)]
    model_lines: List[str] = []
    field_lines: List[str] = []
    for b in bullets:
        if any(k in b for k in ["활용", "분야", "서비스", "정책", "지원"]):
            field_lines.append(b)
        else:
            model_lines.append(b)
    if not model_lines:
        model_lines = ["전지구-지역-극지 연계 모델 체계", "자료동화 및 앙상블 예측 통합", "품질관리 기반 운영 표준화"]
    if not field_lines:
        field_lines = ["국가 정책 수립 지원", "해양·극지 예측 서비스 고도화", "현장 대응 및 국제 협력 활용"]

    model_lines = model_lines[:5]
    field_lines = field_lines[:5]

    for sh in list(slide.shapes):
        name = str(getattr(sh, "name", "") or "")
        if not name.startswith("__OverviewBottom__"):
            continue
        try:
            sh._element.getparent().remove(sh._element)
        except Exception:
            continue

    groups = [
        ("모델 구성", model_lines, int(slide_w * 0.10)),
        ("활용 분야", field_lines, int(slide_w * 0.54)),
    ]
    box_top = int(slide_h * 0.81)
    box_w = int(slide_w * 0.34)
    box_h = int(slide_h * 0.16)
    for title, lines, left in groups:
        shp = slide.shapes.add_textbox(left, box_top, box_w, box_h)
        try:
            shp.name = "__OverviewBottom__group"
        except Exception:
            pass
        shp.fill.background()
        shp.line.fill.background()
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.alignment = PP_ALIGN.LEFT
        _apply_font_all_runs(p0, size_pt=16, bold=True)
        for ln in lines[:5]:
            p = tf.add_paragraph()
            p.text = f"- {ln}"
            p.alignment = PP_ALIGN.LEFT
            _apply_font_all_runs(p, size_pt=13, bold=False)


# 단순 제목 렌더링 함수
def _render_simple_title(
    slide,
    *,
    title: str,
    slide_w: int,
    slide_h: int,
    name_prefix: str,
    top_ratio: float = 0.05,
) -> None:
    for sh in list(slide.shapes):
        nm = str(getattr(sh, "name", "") or "")
        if not nm.startswith(name_prefix):
            continue
        try:
            sh._element.getparent().remove(sh._element)
        except Exception:
            continue
    shp = slide.shapes.add_textbox(
        int(slide_w * 0.06),
        int(slide_h * float(top_ratio)),
        int(slide_w * 0.88),
        int(slide_h * 0.10),
    )
    try:
        shp.name = f"{name_prefix}title"
    except Exception:
        pass
    shp.fill.background()
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = _norm(title)
    _apply_font_all_runs(p, size_pt=30, bold=True)


# 모든 텍스트 도형 제거 함수
def _remove_all_text_shapes(slide) -> None:
    for sh in list(slide.shapes):
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            sh._element.getparent().remove(sh._element)
        except Exception:
            try:
                sh.text_frame.clear()
            except Exception:
                pass


# 생성 이미지 흰 여백 축소 함수
def _tighten_white_margins(image_path: str) -> None:
    """Crop excessive white margins and scale back to original size."""
    try:
        from PIL import Image  # type: ignore
    except Exception:
        return
    try:
        img = Image.open(image_path).convert("RGB")
        w, h = img.size
        px = img.load()
        # non-white pixels
        xs: List[int] = []
        ys: List[int] = []
        for y in range(h):
            for x in range(w):
                r, g, b = px[x, y]
                if (r < 245) or (g < 245) or (b < 245):
                    xs.append(x)
                    ys.append(y)
        if not xs or not ys:
            return
        left = max(0, min(xs))
        right = min(w - 1, max(xs))
        top = max(0, min(ys))
        bottom = min(h - 1, max(ys))
        bw = right - left + 1
        bh = bottom - top + 1
        # skip if already dense enough
        if (bw >= int(w * 0.92)) and (bh >= int(h * 0.92)):
            return
        pad_x = max(2, int(w * 0.02))
        pad_y = max(2, int(h * 0.02))
        left = max(0, left - pad_x)
        right = min(w - 1, right + pad_x)
        top = max(0, top - pad_y)
        bottom = min(h - 1, bottom + pad_y)
        cropped = img.crop((left, top, right + 1, bottom + 1))
        # Resize to original canvas size to maximize visual occupancy.
        fit = cropped.resize((w, h), Image.Resampling.LANCZOS)
        fit.save(image_path)
    except Exception:
        return


# 후처리 단계의 도식/이미지 삽입 함수
def maybe_insert_generated_diagrams(
    pptx_path: str,
    deck_json: Dict[str, Any],
    state: Optional[Dict[str, Any]] = None,
) -> Dict[str, str]:
    if not _enabled(state):
        print("[INFO] Gemini concept images: disabled")
        return {}
    if genai is None or types is None:
        raise RuntimeError("google-genai package is required for Gemini image generation.")

    api_key = os.environ.get("GOOGLE_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("GOOGLE_API_KEY is required for Gemini image generation.")

    preferred = str(
        (state or {}).get("gemini_image_model")
        or os.environ.get("GEMINI_IMAGE_MODEL")
        or "models/gemini-2.5-flash-image"
    ).strip()
    prs = Presentation(pptx_path)
    deck_slides = (deck_json or {}).get("slides") or []
    slide_count = len(prs.slides)
    if slide_count == 0:
        return {}

    targets: List[int] = []
    overview_idx = _find_overview_slide_idx(deck_slides)
    plan_org_idx = _find_plan_orgchart_slide_idx(deck_slides)
    arch_idx = _find_system_arch_slide_idx(deck_slides)
    if overview_idx is not None and 0 <= overview_idx < slide_count:
        targets.append(overview_idx)
    if plan_org_idx is not None and 0 <= plan_org_idx < slide_count and plan_org_idx not in targets:
        targets.append(plan_org_idx)
    if arch_idx is not None and 0 <= arch_idx < slide_count and arch_idx not in targets:
        targets.append(arch_idx)
    targets = targets[:3]
    print(f"[INFO] Gemini 개념 이미지 대상(고정): {targets}")
    if not targets:
        return {}

    client = genai.Client(api_key=api_key)
    model_candidates = _discover_model_candidates(client, preferred)
    print(f"[INFO] Gemini 이미지 모델 후보: {model_candidates[:6]}{'...' if len(model_candidates) > 6 else ''}")

    deck_title = _norm((deck_json or {}).get("deck_title"))
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path((state or {}).get("output_dir") or "output") / "generated_diagrams"

    generated: Dict[str, str] = {}
    inserted_count = 0

    for n, idx in enumerate(targets, 1):
        spec = deck_slides[idx] if idx < len(deck_slides) and isinstance(deck_slides[idx], dict) else {}
        section = _norm(spec.get("section"))
        title = _norm(spec.get("slide_title"))
        prompt_type = _norm(spec.get("image_prompt_type")).lower()
        if prompt_type == "plan_orgchart_fixed":
            slide = prs.slides[idx]
            fixed_bg = _resolve_plan_orgchart_bg_image_path("", state=state)
            _remove_all_text_shapes(slide)
            _render_simple_title(
                slide,
                title=_norm(spec.get("slide_title") or "추진 계획 조직도"),
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
                name_prefix="__PlanOrg__",
            )
            _apply_basic_origin_background(slide, int(prs.slide_width), int(prs.slide_height), state=state)
            slot = _basic_main_panel_slot(int(prs.slide_width), int(prs.slide_height))
            _remove_overlapping_shapes(
                slide,
                slot,
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
                margin_ratio=0.0,
                preserve_text_shapes=True,
            )
            ok = _insert_picture_front(slide, fixed_bg, slot)
            if ok:
                _overlay_plan_orgchart_texts(slide, slot)
                generated[f"slide_{idx}_image_path"] = fixed_bg
                generated[f"slide_{idx}_overlay"] = "plan_orgchart_fixed"
                inserted_count += 1
                spec["layout"] = "text_image"
                spec["visual_slot"] = "panel_main"
                spec["image_needed"] = True
                spec["image_type"] = "diagram"
                continue
            spec["layout"] = "text_only"
            spec["image_needed"] = False
            spec["image_type"] = "none"
            continue

        if prompt_type == "system_architecture":
            # 시스템 아키텍처는 새 이미지를 생성하지 않음
            # 사전 준비 배경 이미지와 좌표 기반 라벨 오버레이만 사용
            slide = prs.slides[idx]
            fixed_bg = _resolve_arch_bg_image_path("", state=state)
            overlay_ok = add_architecture_overlay(
                slide,
                fixed_bg,
                state=state,
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
            )
            if overlay_ok:
                generated[f"slide_{idx}_image_path"] = fixed_bg
                generated[f"slide_{idx}_overlay"] = "system_architecture"
                inserted_count += 1
                spec["layout"] = "text_image"
                spec["visual_slot"] = "panel_main"
                spec["image_needed"] = True
                spec["image_type"] = "diagram"
                continue
            spec["layout"] = "text_only"
            spec["image_needed"] = False
            spec["image_type"] = "none"
            continue

        prompt = _build_prompt(deck_title, section, title, prompt_type=prompt_type)
        out_path = out_dir / f"concept_{n}_{ts}.png"
        img_path = _generate_one_image(
            client,
            model_candidates,
            prompt,
            out_path,
            max_retries=int((state or {}).get("gemini_image_retry_count") or os.environ.get("GEMINI_IMAGE_RETRY_COUNT") or 1),
        )
        if not img_path:
            spec["layout"] = "text_only"
            spec["image_needed"] = False
            spec["image_type"] = "none"
            continue
        if prompt_type == "overview_last":
            _tighten_white_margins(img_path)

        if prompt_type == "plan":
            bullets = [str(b or "").strip() for b in (spec.get("bullets") or []) if str(b or "").strip()]
            line1 = _norm(spec.get("key_message") or "")
            line2 = bullets[0] if bullets else ""
            _render_plan_text_header(
                prs.slides[idx],
                title=_norm(spec.get("slide_title") or "추진 계획"),
                line1=line1,
                line2=line2,
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
            )
            slot = _plan_image_slot(int(prs.slide_width), int(prs.slide_height))
            spec["layout"] = "text_image"
            spec["visual_slot"] = "right_large"
            spec["image_needed"] = True
            spec["image_type"] = "diagram"
        elif prompt_type == "overview_last":
            slide = prs.slides[idx]
            _remove_all_text_shapes(slide)
            _render_simple_title(
                slide,
                title=_norm(spec.get("slide_title") or "연구개요"),
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
                name_prefix="__Overview__",
                top_ratio=0.065,
            )
            _apply_basic_origin_background(slide, int(prs.slide_width), int(prs.slide_height), state=state)
            slot = _overview_center_image_slot(int(prs.slide_width), int(prs.slide_height))
            _remove_overlapping_shapes(
                slide,
                slot,
                slide_w=int(prs.slide_width),
                slide_h=int(prs.slide_height),
                margin_ratio=0.10,
                preserve_text_shapes=True,
            )
            if _insert_picture_front(slide, img_path, slot):
                _overview_bottom_text_groups(
                    slide,
                    spec,
                    slide_w=int(prs.slide_width),
                    slide_h=int(prs.slide_height),
                )
                generated[f"slide_{idx}_image_path"] = img_path
                inserted_count += 1
                spec["layout"] = "text_image"
                spec["visual_slot"] = "panel_main"
                spec["image_needed"] = True
                spec["image_type"] = "diagram"
                continue
            spec["layout"] = "text_only"
            spec["image_needed"] = False
            spec["image_type"] = "none"
            continue
        else:
            slot = _text_image_slot(int(prs.slide_width), int(prs.slide_height))
            spec["layout"] = "text_image"
            spec["visual_slot"] = "right_large"
            spec["image_needed"] = True
            spec["image_type"] = "diagram"

        _remove_overlapping_shapes(
            prs.slides[idx],
            slot,
            slide_w=int(prs.slide_width),
            slide_h=int(prs.slide_height),
            margin_ratio=0.06,
            preserve_text_shapes=False,
        )
        ok = _insert_picture(prs.slides[idx], img_path, slot)
        if ok:
            generated[f"slide_{idx}_image_path"] = img_path
            inserted_count += 1
        else:
            spec["layout"] = "text_only"
            spec["image_needed"] = False
            spec["image_type"] = "none"

    prs.save(pptx_path)
    print(f"[INFO] Gemini 개념 이미지 삽입 결과: inserted={inserted_count}")
    return generated


# 슬라이드 텍스트 추출 함수
def _slide_text(slide) -> str:
    out: List[str] = []
    for sh in list(slide.shapes):
        if not getattr(sh, "has_text_frame", False):
            continue
        t = _norm(getattr(sh, "text", ""))
        if t:
            out.append(t)
    return " ".join(out)


# 개발용 로컬 확인 함수: 시스템 아키텍처 오버레이만 적용
def preview_architecture_overlay_only(
    input_pptx: str,
    output_pptx: Optional[str] = None,
    *,
    slide_index: Optional[int] = None,
    title_keyword: str = "시스템 아키텍처",
    bg_image_path: str = "features/ppt_maker/background/system.png",
) -> Dict[str, Any]:
    prs = Presentation(input_pptx)
    targets: List[int] = []

    if slide_index is not None:
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"slide_index out of range: {slide_index}")
        targets = [slide_index]
    else:
        kw = _norm(title_keyword)
        for i, s in enumerate(prs.slides):
            if kw and kw in _slide_text(s):
                targets.append(i)

    if not targets:
        raise RuntimeError("시스템 아키텍처 오버레이 확인 대상 슬라이드를 찾지 못했습니다.")

    state = {"system_arch_bg_image_path": bg_image_path}
    applied = 0
    for idx in targets:
        if add_architecture_overlay(
            prs.slides[idx],
            bg_image_path,
            state=state,
            slide_w=int(prs.slide_width),
            slide_h=int(prs.slide_height),
        ):
            applied += 1

    if output_pptx:
        out_path = output_pptx
    else:
        p2 = Path(input_pptx)
        out_path = str(p2.with_name(f"{p2.stem}_arch_check{p2.suffix}"))

    prs.save(out_path)
    return {"output_pptx": out_path, "targets": targets, "applied": applied}


# 개발용 로컬 확인 함수: 추진 계획 조직도 오버레이만 적용
def preview_plan_orgchart_overlay_only(
    input_pptx: str,
    output_pptx: Optional[str] = None,
    *,
    slide_index: Optional[int] = None,
    title_keyword: str = "추진 계획 조직도",
    bg_image_path: str = "features/ppt_maker/background/Hierarchical organization diagram..png",
    line1: str = "",
    line2: str = "",
) -> Dict[str, Any]:
    prs = Presentation(input_pptx)
    targets: List[int] = []

    if slide_index is not None:
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(f"slide_index out of range: {slide_index}")
        targets = [slide_index]
    else:
        kw = _norm(title_keyword)
        for i, s in enumerate(prs.slides):
            if kw and kw in _slide_text(s):
                targets.append(i)

    if not targets:
        raise RuntimeError("추진 계획 조직도 오버레이 확인 대상 슬라이드를 찾지 못했습니다.")

    state = {"plan_orgchart_bg_image_path": bg_image_path}
    applied = 0
    for idx in targets:
        slide = prs.slides[idx]
        fixed_bg = _resolve_plan_orgchart_bg_image_path("", state=state)
        _apply_basic_origin_background(slide, int(prs.slide_width), int(prs.slide_height), state=state)
        slot = _basic_main_panel_slot(int(prs.slide_width), int(prs.slide_height))
        _remove_overlapping_shapes(
            slide,
            slot,
            slide_w=int(prs.slide_width),
            slide_h=int(prs.slide_height),
            margin_ratio=0.0,
            preserve_text_shapes=False,
        )
        if _insert_picture_front(slide, fixed_bg, slot):
            _overlay_plan_orgchart_texts(slide, slot)
            applied += 1

    if output_pptx:
        out_path = output_pptx
    else:
        p2 = Path(input_pptx)
        out_path = str(p2.with_name(f"{p2.stem}_plan_orgchart_check{p2.suffix}"))

    prs.save(out_path)
    return {"output_pptx": out_path, "targets": targets, "applied": applied}


# 개발용 로컬 확인 함수: 시스템 아키텍처와 추진 계획 조직도만 남긴 미리보기 생성
def preview_arch_and_plan_only(
    input_pptx: str,
    output_pptx: str,
    *,
    arch_slide_index: int,
    plan_slide_index: int,
    arch_bg_image_path: str = "features/ppt_maker/background/system.png",
    plan_bg_image_path: str = "features/ppt_maker/background/Hierarchical organization diagram..png",
) -> Dict[str, Any]:
    prs = Presentation(input_pptx)

    if not (0 <= arch_slide_index < len(prs.slides)):
        raise IndexError(f"arch_slide_index out of range: {arch_slide_index}")
    if not (0 <= plan_slide_index < len(prs.slides)):
        raise IndexError(f"plan_slide_index out of range: {plan_slide_index}")

    add_architecture_overlay(
        prs.slides[arch_slide_index],
        arch_bg_image_path,
        state={"system_arch_bg_image_path": arch_bg_image_path},
        slide_w=int(prs.slide_width),
        slide_h=int(prs.slide_height),
    )

    plan_slide = prs.slides[plan_slide_index]
    _apply_basic_origin_background(plan_slide, int(prs.slide_width), int(prs.slide_height), state={})
    plan_slot = _basic_main_panel_slot(int(prs.slide_width), int(prs.slide_height))
    _remove_overlapping_shapes(
        plan_slide,
        plan_slot,
        slide_w=int(prs.slide_width),
        slide_h=int(prs.slide_height),
        margin_ratio=0.0,
        preserve_text_shapes=False,
    )
    _insert_picture_front(
        plan_slide,
        _resolve_plan_orgchart_bg_image_path("", state={"plan_orgchart_bg_image_path": plan_bg_image_path}),
        plan_slot,
    )
    _overlay_plan_orgchart_texts(plan_slide, plan_slot)

    keep = sorted({arch_slide_index, plan_slide_index})
    for i in reversed(range(len(prs.slides))):
        if i in keep:
            continue
        sld_id_lst = prs.slides._sldIdLst  # pylint: disable=protected-access
        sld = sld_id_lst[i]
        rel_id = sld.rId
        sld_id_lst.remove(sld)
        prs.part.drop_rel(rel_id)

    prs.save(output_pptx)
    return {"output_pptx": output_pptx, "kept_indices": keep}


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="시스템 아키텍처 오버레이 로컬 확인(Gamma API 미사용)")
    parser.add_argument("--input", required=True, help="입력 PPTX 경로")
    parser.add_argument("--output", default="", help="출력 PPTX 경로")
    parser.add_argument("--slide-index", type=int, default=None, help="0부터 시작하는 대상 슬라이드 번호")
    parser.add_argument("--keyword", default="시스템 아키텍처", help="slide-index 생략 시 대상 탐색용 제목 키워드")
    parser.add_argument("--bg", default="features/ppt_maker/background/system.png", help="배경 이미지 경로")
    args = parser.parse_args()

    result = preview_architecture_overlay_only(
        input_pptx=args.input,
        output_pptx=(args.output or None),
        slide_index=args.slide_index,
        title_keyword=args.keyword,
        bg_image_path=args.bg,
    )
    print(f"[INFO] 시스템 아키텍처 오버레이 확인 파일 저장: {result}")
